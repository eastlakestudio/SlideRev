import os
import cv2
import numpy as np
import fitz  # PyMuPDF
from paddleocr import PaddleOCR
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import logging

# 配置日誌
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')
logger = logging.getLogger(__name__)

class SlideReverseCore:
    def __init__(self, lang='ch', use_gpu=False):
        """
        初始化 PaddleOCR 和相關設定
        """
        self.ocr = PaddleOCR(use_angle_cls=True, lang=lang, use_gpu=use_gpu, show_log=False)
        self.output_resolution = 300  # DPI

    def process_pdf(self, pdf_path, output_pptx_path, progress_callback=None):
        """
        主處理 Pipeline
        """
        try:
            logger.info(f"正在讀取 PDF: {pdf_path}")
            doc = fitz.open(pdf_path)
            total_pages = len(doc)
            
            prs = Presentation()
            # 設置 16:9 比例 (13.33 x 7.5 英吋)
            prs.slide_width = Inches(13.33)
            prs.slide_height = Inches(7.5)

            temp_dir = "temp_slides"
            os.makedirs(temp_dir, exist_ok=True)

            for i in range(total_pages):
                logger.info(f"處理第 {i+1}/{total_pages} 頁")
                if progress_callback:
                    progress_callback(i + 1, total_pages, "正在處理頁面...")

                page = doc.load_page(i)
                # 設置渲染比例 (提高解析度)
                pix = page.get_pixmap(matrix=fitz.Matrix(300/72, 300/72))
                
                # 將 pixmap 轉換為 numpy array (OpenCV 格式)
                img_cv = np.frombuffer(pix.samples, dtype=np.uint8).reshape(pix.h, pix.w, pix.n)
                if pix.n == 4: # RGBA -> BGR
                    img_cv = cv2.cvtColor(img_cv, cv2.COLOR_RGBA2BGR)
                else: # RGB -> BGR
                    img_cv = cv2.cvtColor(img_cv, cv2.COLOR_RGB2BGR)
                
                # 1. OCR 識別
                ocr_results = self.ocr.ocr(img_cv, cls=True)[0]
                
                # 2. 生成 Mask 並進行 Inpainting
                mask = np.zeros(img_cv.shape[:2], dtype=np.uint8)
                text_info = []

                if ocr_results:
                    for line in ocr_results:
                        box = line[0]  # [[x,y], [x,y], [x,y], [x,y]]
                        text = line[1][0]
                        conf = line[1][1]
                        
                        # 記錄座標與內容
                        text_info.append({'box': box, 'text': text})
                        
                        # 在 Mask 上標註文字區域
                        pts = np.array(box, np.int32).reshape((-1, 1, 2))
                        cv2.fillPoly(mask, [pts], 255)

                # 3. 修復背景
                # 使用 INPAINT_TELEA 或 INPAINT_NS
                inpainted_bg = cv2.inpaint(img_cv, mask, 3, cv2.INPAINT_TELEA)
                
                # 儲存修復後的背景圖
                bg_path = os.path.join(temp_dir, f"bg_page_{i+1}.jpg")
                cv2.imwrite(bg_path, inpainted_bg)

                # 4. 建立投影片
                slide_layout = prs.slide_layouts[6]  # 空白佈局
                slide = prs.slides.add_slide(slide_layout)
                
                # 將修復後的背景圖鋪滿投影片
                slide.shapes.add_picture(bg_path, 0, 0, width=prs.slide_width, height=prs.slide_height)

                # 5. 插入文字方塊
                img_h, img_w = img_cv.shape[:2]
                for item in text_info:
                    box = item['box']
                    text = item['text']
                    
                    # 計算文字方塊的邊界框 (Bounding Box)
                    x_min = min([p[0] for p in box])
                    y_min = min([p[1] for p in box])
                    x_max = max([p[0] for p in box])
                    y_max = max([p[1] for p in box])
                    
                    w = x_max - x_min
                    h = y_max - y_min
                    
                    # 映射到 PPT 座標 (英吋)
                    left = Inches(x_min / img_w * 13.33)
                    top = Inches(y_min / img_h * 7.5)
                    width = Inches(w / img_w * 13.33)
                    height = Inches(h / img_h * 7.5)

                    txBox = slide.shapes.add_textbox(left, top, width, height)
                    tf = txBox.text_frame
                    tf.text = text
                    
                    # 嘗試取樣顏色 (取 box 中心點的顏色)
                    center_x, center_y = int((x_min + x_max)/2), int((y_min + y_max)/2)
                    if 0 <= center_x < img_w and 0 <= center_y < img_h:
                        color = img_cv[center_y, center_x]  # BGR
                        run = tf.paragraphs[0].runs[0]
                        run.font.color.rgb = RGBColor(color[2], color[1], color[0])
                        
                    # 設置字體大小 (根據 box 高度估計)
                    font_size = h / img_h * 7.5 * 72 # 1 inch = 72 points
                    tf.paragraphs[0].runs[0].font.size = Pt(max(6, font_size * 0.8)) # 稍微縮小一點避免溢出

            # 儲存結果
            prs.save(output_pptx_path)
            logger.info(f"轉換完成: {output_pptx_path}")
            
            # 清理臨時文件
            for f in os.listdir(temp_dir):
                os.remove(os.path.join(temp_dir, f))
            os.rmdir(temp_dir)
            
            return True

        except Exception as e:
            logger.error(f"轉換過程中發生錯誤: {str(e)}")
            return False

if __name__ == "__main__":
    # 簡單測試代碼
    # core = SlideReverseCore()
    # core.process_pdf("test.pdf", "output.pptx")
    pass
