import json
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE

def hex_to_rgb(hex_str):
    hex_str = hex_str.lstrip('#')
    return tuple(int(hex_str[i:i+2], 16) for i in (0, 2, 4))

def invert_rgb(rgb):
    return (255 - rgb[0], 255 - rgb[1], 255 - rgb[2])

def create_pptx(json_path, slides_dir, output_path):
    with open(json_path, 'r', encoding='utf-8') as f:
        data = json.load(f)
    
    prs = Presentation()
    
    # User noted font is too large, so we'll start with 1.0 scale 
    # and let 'TEXT_TO_FIT_SHAPE' shrink it if it's too big for the box.
    font_scale = 1.0 # Standard scale, let auto-size handle it.
    
    for page in data:
        # Set slide size based on first page
        if page['pageIndex'] == 0:
            prs.slide_width = Inches(page['width'] / 72.0)
            prs.slide_height = Inches(page['height'] / 72.0)
            
        slide = prs.slides.add_slide(prs.slide_layouts[6]) # blank layout
        
        # 1. Add Background Image (Original PDF page)
        img_path = os.path.join(slides_dir, f"page_{page['pageIndex'] + 1}.jpg")
        if os.path.exists(img_path):
            slide.shapes.add_picture(img_path, 0, 0, width=prs.slide_width, height=prs.slide_height)
        
        # 2. Add Spans (Rectangles + Text)
        for span in page['spans']:
            left = Inches(span['x'] / 72.0)
            top = Inches(span['y'] / 72.0)
            width = Inches(span['width'] / 72.0)
            height = Inches(span['height'] / 72.0)
            
            # Draw Background Box (Opaque)
            text_color_rgb = hex_to_rgb(span['color'])
            bg_color_rgb = invert_rgb(text_color_rgb)
            
            # Create a shape that acts as BOTH background box and text box
            shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, left, top, width, height
            )
            
            # Fill color (opaque)
            shape.fill.solid()
            shape.fill.foreground_color.rgb = RGBColor(*bg_color_rgb)
            shape.line.fill.background() # No border
            
            # Add Text Frame
            tf = shape.text_frame
            tf.text = span['text']
            tf.word_wrap = True # Ensure it wraps if needed
            
            # AUTO-RESIZE: Text to fit shape (if font too large, it shrinks)
            tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            
            # Padding/Margins (Set to 0 as requested for alignment precision)
            tf.margin_bottom = 0
            tf.margin_top = 0
            tf.margin_left = 0
            tf.margin_right = 0
            
            # Text Formatting
            p = tf.paragraphs[0]
            # p.font.name = 'PingFang SC' # Apply to paragraph/run
            
            for run in p.runs:
                run.font.size = Pt(span['fontSize'] * font_scale)
                run.font.color.rgb = RGBColor(*text_color_rgb)
                run.font.name = 'PingFang SC'
            
    prs.save(output_path)
    print(f"✅ Success! PPTX with Opaque Backgrounds & Auto-fit saved to {output_path}")

if __name__ == "__main__":
    create_pptx("test_origin_extracted.json", "temp_slides", "test_origin_final.pptx")
