import sys
import json
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def create_pptx(data_path, output_path):
    with open(data_path, 'r') as f:
        data = json.load(f)
    
    prs = Presentation()
    
    # Use the first image's aspect ratio to set slide size
    if data and len(data) > 0:
        first_page = data[0]
        # In PDF/Points (1/72 inch). Standard is 720x540 for 4:3.
        # We try to match the image pixel size if provided, or use defaults.
        img_w = first_page.get('image_width', 800)
        img_h = first_page.get('image_height', 600)
        prs.slide_width = Inches(img_w / 72.0)
        prs.slide_height = Inches(img_h / 72.0)

    for page in data:
        slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)
        
        # 1. Add Background (Inpainted Image)
        img_path = page.get('background_path')
        if img_path and os.path.exists(img_path):
            slide.shapes.add_picture(img_path, 0, 0, width=prs.slide_width, height=prs.slide_height)
        
        # 2. Add Text Items
        for item in page.get('items', []):
            if not item.get('is_erased', False): continue
            
            text = item.get('text', "")
            if not text: continue
            
            # Coordinates are normalized (0-1) relative to slide size
            left = prs.slide_width * item['x']
            top = prs.slide_height * (1 - item['y'] - item['h']) # Flip Y for PPTX (top-origin)
            width = prs.slide_width * item['w']
            height = prs.slide_height * item['h']
            
            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            tf.word_wrap = True
            
            p = tf.paragraphs[0]
            p.text = text
            p.alignment = PP_ALIGN.LEFT
            
            # Styling
            run = p.runs[0]
            run.font.size = Pt(item.get('font_size', 12) * (72.0 / 300.0)) # Convert pixels back to points (assuming 300dpi)
            
            # Color (Hex to RGB)
            color_hex = item.get('color', "#000000").lstrip('#')
            if len(color_hex) == 6:
                r, g, b = tuple(int(color_hex[i:i+2], 16) for i in (0, 2, 4))
                run.font.color.rgb = RGBColor(r, g, b)
            
            # Rotation
            txBox.rotation = -item.get('rotation', 0) # PPTX uses clockwise rotation

    prs.save(output_path)
    print(f"✅ Successfully exported to {output_path}")

if __name__ == "__main__":
    if len(sys.argv) < 3:
        print("Usage: python3 export_to_pptx.py <json_data> <output_pptx>")
    else:
        create_pptx(sys.argv[1], sys.argv[2])
